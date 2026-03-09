"""
generate_report.py
==================
Monarch Analytics — Automated PowerPoint Report Generator

Generates a 20-slide branded audit report using python-pptx.
Charts from analysis are automatically embedded into slides.
"""

import os
import json
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand Colors ──────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)
NAVY        = RGBColor(0x1B, 0x4F, 0x72)
BLUE        = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_BLUE  = RGBColor(0xA8, 0xDA, 0xDC)
ORANGE      = RGBColor(0xF4, 0xA2, 0x61)
RED         = RGBColor(0xE6, 0x39, 0x46)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

VIZ_DIR     = "audit_workspace/visualizations"
ANALYSIS_DIR = "audit_workspace/analysis"
OUT_PATH    = "audit_workspace/report_outputs/business_audit_report.pptx"


# ── Helper Functions ──────────────────────────────────────────────────────────

def load_kpis() -> dict:
    path = os.path.join(ANALYSIS_DIR, "kpi_results.json")
    if os.path.isfile(path):
        with open(path) as f:
            return json.load(f)
    return {}


def load_leaks() -> dict:
    path = os.path.join(ANALYSIS_DIR, "revenue_leak_report.json")
    if os.path.isfile(path):
        with open(path) as f:
            return json.load(f)
    return {}


def set_slide_bg(slide, color: RGBColor):
    """Set a solid background color for a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_image(slide, img_path, left, top, width, height=None):
    """Insert an image into a slide if it exists."""
    if os.path.isfile(img_path):
        if height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width))
    else:
        add_textbox(slide, f"[Chart not found: {os.path.basename(img_path)}]",
                    left, top, width, height or 2, font_size=10, color=ORANGE)


def add_divider(slide, top, color=BLUE):
    """Add a horizontal line divider."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(0.4), Inches(top), Inches(9.6), Inches(top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1)


def add_section_header(slide, title, subtitle=""):
    """Standard slide header layout."""
    add_textbox(slide, title, 0.4, 0.2, 9.2, 0.7,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 0.85, 9.2, 0.4,
                    font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    add_divider(slide, 1.2)


def add_kpi_box(slide, label, value, left, top, w=2.1, h=1.2,
                bg_color=NAVY, val_color=ORANGE):
    """Add a styled KPI metric card."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = BLUE

    tb = shape.text_frame
    tb.word_wrap = True
    tb.text = value
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = val_color
    run.font.name = "Calibri"

    p2 = tb.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = LIGHT_BLUE
    r2.font.name = "Calibri"


# ── Slide Builders ────────────────────────────────────────────────────────────

def slide_title(prs, kpis):
    """Slide 1: Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Monarch branding bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(slide, "MONARCH ANALYTICS", 0.4, 0.15, 9.2, 0.7,
                font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

    add_textbox(slide, "Business Performance\nAudit Report", 0.8, 1.6, 8.4, 2.0,
                font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, "Confidential | Prepared by Monarch Analytics",
                0.8, 3.8, 8.4, 0.5, font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

    date_str = datetime.now().strftime("%B %Y")
    add_textbox(slide, date_str, 0.8, 4.3, 8.4, 0.4,
                font_size=13, color=ORANGE, align=PP_ALIGN.LEFT)

    # Decorative accent
    accent = slide.shapes.add_shape(1, Inches(0), Inches(6.2), Inches(10), Inches(0.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    add_textbox(slide, "DATA-DRIVEN DECISIONS  |  REVENUE GROWTH  |  BUSINESS CLARITY",
                0.4, 6.25, 9.2, 0.5, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_executive_summary(prs, kpis, leaks):
    """Slide 2: Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Executive Summary",
                       "Key findings from the business performance audit")

    s = kpis.get("sales", {})
    m = kpis.get("marketing", {})
    c = kpis.get("customer", {})
    f = kpis.get("financial", {})
    n_leaks = leaks.get("total_leaks_detected", 0)
    impact = leaks.get("total_estimated_revenue_impact", 0)

    summary_lines = [
        f"• Total Revenue Analyzed: ₹{s.get('total_revenue', 0):,.0f}",
        f"• Lead-to-Sale Conversion: {s.get('lead_conversion_rate', 0)*100:.1f}%",
        f"• Customer Acquisition Cost: ₹{m.get('customer_acquisition_cost', 0):,.0f}",
        f"• Return on Ad Spend: {m.get('return_on_ad_spend', 0):.1f}x",
        f"• Customer LTV: ₹{c.get('customer_lifetime_value', 0):,.0f}",
        f"• Customer Retention Rate: {c.get('retention_rate', 0)*100:.1f}%",
        f"• Gross Margin: {f.get('gross_margin', 0)*100:.1f}%",
        f"• Revenue Leaks Detected: {n_leaks}",
        f"• Estimated Revenue at Risk: ₹{impact:,.0f}",
    ]

    add_textbox(slide, "\n".join(summary_lines), 0.4, 1.4, 5.5, 5.0,
                font_size=14, color=WHITE)

    # Leaks box
    box = slide.shapes.add_shape(1, Inches(6.2), Inches(1.4), Inches(3.4), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.color.rgb = RED

    tb = box.text_frame
    tb.text = f"{n_leaks}"
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(52)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RED

    p2 = tb.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Revenue Leaks\nDetected"
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE

    # Impact box
    box2 = slide.shapes.add_shape(1, Inches(6.2), Inches(3.8), Inches(3.4), Inches(1.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = NAVY
    box2.line.color.rgb = ORANGE

    tb2 = box2.text_frame
    tb2.text = f"₹{impact:,.0f}"
    p3 = tb2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.runs[0].font.size = Pt(22)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = ORANGE

    p4 = tb2.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Revenue at Risk"
    r4.font.size = Pt(12)
    r4.font.color.rgb = LIGHT_BLUE


def slide_key_metrics(prs, kpis):
    """Slide 3: KPI Snapshot"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Key Metrics Snapshot", "Core KPIs across sales, marketing, and customer")

    s = kpis.get("sales", {})
    m = kpis.get("marketing", {})
    c = kpis.get("customer", {})
    f = kpis.get("financial", {})

    metrics = [
        ("Total Revenue", f"₹{s.get('total_revenue', 0):,.0f}"),
        ("Avg Deal Value", f"₹{s.get('average_deal_value', 0):,.0f}"),
        ("Win Rate", f"{s.get('lead_conversion_rate', 0)*100:.1f}%"),
        ("ROAS", f"{m.get('return_on_ad_spend', 0):.1f}x"),
        ("CAC", f"₹{m.get('customer_acquisition_cost', 0):,.0f}"),
        ("LTV", f"₹{c.get('customer_lifetime_value', 0):,.0f}"),
        ("Retention", f"{c.get('retention_rate', 0)*100:.1f}%"),
        ("Gross Margin", f"{f.get('gross_margin', 0)*100:.1f}%"),
    ]

    cols = 4
    for i, (label, val) in enumerate(metrics):
        col = i % cols
        row = i // cols
        left = 0.3 + col * 2.35
        top = 1.5 + row * 1.5
        add_kpi_box(slide, label, val, left, top)


def slide_revenue_chart(prs):
    """Slide 4: Revenue Trend"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Revenue Performance", "Monthly revenue trend analysis")
    add_image(slide, os.path.join(VIZ_DIR, "revenue_trend.png"), 0.5, 1.4, 9.0, 5.0)


def slide_sales_funnel(prs):
    """Slide 5: Sales Funnel"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Sales Funnel Analysis", "Lead-to-close conversion performance")
    add_image(slide, os.path.join(VIZ_DIR, "sales_funnel_chart.png"), 0.5, 1.4, 9.0, 5.0)


def slide_marketing(prs):
    """Slide 6: Marketing Performance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Marketing Efficiency", "Channel performance and spend effectiveness")
    add_image(slide, os.path.join(VIZ_DIR, "channel_performance.png"), 0.3, 1.4, 9.4, 5.0)


def slide_cac_ltv(prs):
    """Slide 7: CAC vs LTV"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "CAC vs Customer Lifetime Value",
                       "Sustainability of customer acquisition investment")
    add_image(slide, os.path.join(VIZ_DIR, "cac_vs_ltv.png"), 1.0, 1.4, 8.0, 4.8)


def slide_marketing_funnel(prs):
    """Slide 8: Marketing Funnel"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Marketing Conversion Funnel",
                       "Impression-to-conversion drop-off analysis")
    add_image(slide, os.path.join(VIZ_DIR, "marketing_funnel.png"), 0.5, 1.4, 9.0, 5.0)


def slide_customer_ltv(prs):
    """Slide 9: LTV Distribution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Customer Lifetime Value Distribution",
                       "Revenue concentration and value segmentation")
    add_image(slide, os.path.join(VIZ_DIR, "ltv_distribution.png"), 0.5, 1.4, 9.0, 5.0)


def slide_customer_segments(prs):
    """Slide 10: Customer Segments"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Customer Segment Analysis",
                       "Revenue and retention by customer tier")
    add_image(slide, os.path.join(VIZ_DIR, "customer_cohort_retention.png"), 0.3, 1.4, 9.4, 5.0)


def slide_pareto(prs, drivers):
    """Slide 11: Pareto Revenue"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    pct = drivers.get("pareto_top_20pct_revenue_share", 0) if drivers else 0
    add_section_header(slide, "Pareto Revenue Analysis",
                       f"Top 20% of customers drive {pct}% of total revenue")
    add_image(slide, os.path.join(VIZ_DIR, "pareto_revenue_chart.png"), 0.3, 1.4, 9.4, 5.0)


def slide_revenue_leaks(prs, leaks):
    """Slide 12: Revenue Leak Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Revenue Leak Detection Results",
                       "Automated analysis of business performance gaps")

    leak_list = leaks.get("leaks", [])
    if not leak_list:
        add_textbox(slide, "✅ No significant revenue leaks detected.",
                    0.5, 2.0, 9.0, 1.0, font_size=18, color=GREEN)
        return

    for i, leak in enumerate(leak_list[:5]):
        top = 1.5 + i * 1.0
        score = leak.get("impact_score", 0)
        color = RED if score >= 8 else ORANGE if score >= 6 else LIGHT_BLUE

        box = slide.shapes.add_shape(1, Inches(0.3), Inches(top), Inches(9.4), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.color.rgb = color

        tb = box.text_frame
        tb.word_wrap = True
        tb.text = f"[{score}/10] {leak.get('flag', '')}  —  {leak.get('description', '')}"
        p = tb.paragraphs[0]
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = "Calibri"


def slide_fix_improve_ignore(prs, leaks):
    """Slide 13: Fix / Improve / Ignore Framework"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Fix | Improve | Ignore Framework",
                       "Prioritized action plan based on impact and effort")

    leak_list = leaks.get("leaks", [])
    fix_items    = [l for l in leak_list if l.get("impact_score", 0) >= 8]
    improve_items = [l for l in leak_list if 5 <= l.get("impact_score", 0) < 8]
    ignore_items  = [l for l in leak_list if l.get("impact_score", 0) < 5]

    sections = [
        ("🔴 FIX (Immediate)", fix_items, RED, 0.3),
        ("🟡 IMPROVE (This Quarter)", improve_items, ORANGE, 3.5),
        ("🟢 IGNORE (Low Impact)", ignore_items, GREEN, 6.7),
    ]

    for label, items, color, left in sections:
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.4), Inches(3.0), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.color.rgb = color

        tb = box.text_frame
        tb.word_wrap = True
        tb.text = label
        p = tb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = color

        for item in items[:4]:
            p2 = tb.add_paragraph()
            r = p2.add_run()
            r.text = f"• {item.get('rule', '').replace('_', ' ').title()}"
            r.font.size = Pt(10)
            r.font.color.rgb = WHITE

        if not items:
            p2 = tb.add_paragraph()
            r = p2.add_run()
            r.text = "• None identified"
            r.font.size = Pt(10)
            r.font.color.rgb = LIGHT_BLUE


def slide_recommendations(prs, leaks):
    """Slide 14: Top Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Recommended Actions",
                       "Prioritized steps to recover revenue and improve performance")

    leak_list = leaks.get("leaks", [])
    top_leaks = sorted(leak_list, key=lambda x: x.get("impact_score", 0), reverse=True)[:5]

    for i, leak in enumerate(top_leaks):
        top = 1.5 + i * 1.0
        add_textbox(slide,
                    f"{i+1}.  {leak.get('recommended_action', '')}",
                    0.5, top, 9.0, 0.85, font_size=12, color=WHITE)


def slide_roadmap(prs):
    """Slide 15: Implementation Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "90-Day Implementation Roadmap",
                       "Execution timeline for audit recommendations")

    phases = [
        ("Days 1–30\nQuick Wins", "• Fix top funnel leaks\n• Pause low-ROAS ad sets\n• Launch churn survey", ORANGE),
        ("Days 31–60\nGrowth Sprint", "• Implement retention program\n• Pricing review\n• Sales enablement rollout", BLUE),
        ("Days 61–90\nScale Up", "• Review KPI dashboard\n• Measure improvements\n• Expand winning channels", GREEN),
    ]

    for i, (title, bullets, color) in enumerate(phases):
        left = 0.3 + i * 3.3
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.5), Inches(3.1), Inches(4.6))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.color.rgb = color

        tb = box.text_frame
        tb.word_wrap = True
        tb.text = title
        p = tb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = color

        for line in bullets.split("\n"):
            p2 = tb.add_paragraph()
            r = p2.add_run()
            r.text = line
            r.font.size = Pt(11)
            r.font.color.rgb = WHITE


def slide_financial_impact(prs, leaks, kpis):
    """Slide 16: Financial Impact Estimate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Financial Impact Estimate",
                       "Revenue recovery potential from addressing identified leaks")

    total_impact = leaks.get("total_estimated_revenue_impact", 0)
    total_rev = kpis.get("sales", {}).get("total_revenue", 1)
    pct = (total_impact / total_rev * 100) if total_rev else 0

    lines = [
        f"Total Revenue Analyzed:          ₹{total_rev:,.0f}",
        f"Revenue at Risk (Leaks):          ₹{total_impact:,.0f}",
        f"Revenue at Risk (%):               {pct:.1f}%",
        "",
        "If all recommendations are implemented:",
        f"  Estimated Revenue Recovery:  ₹{total_impact * 0.6:,.0f} – ₹{total_impact * 0.9:,.0f}",
        f"  Conservative ROI on Audit:    {max(3, total_impact * 0.5 / 50000):.0f}x",
    ]

    add_textbox(slide, "\n".join(lines), 0.5, 1.5, 9.0, 4.5,
                font_size=15, color=WHITE)


def slide_retainer_offer(prs):
    """Slide 17: Monarch Retainer Offer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(slide, "MONARCH ANALYTICS", 0.4, 0.12, 4, 0.6,
                font_size=13, bold=True, color=ORANGE)

    add_textbox(slide, "Continue the Growth Journey", 0.8, 1.2, 8.4, 0.8,
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Monarch Analytics Retainer — Monthly Intelligence Partnership",
                0.8, 2.0, 8.4, 0.5, font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    offerings = [
        "✅  Monthly KPI Dashboard Updates",
        "✅  Quarterly Business Performance Audits",
        "✅  Dedicated Analytics Support (Slack / Email)",
        "✅  Revenue Leak Monitoring Alerts",
        "✅  Custom Dashboard & Reporting",
        "✅  Strategic Growth Advisory Sessions",
    ]
    add_textbox(slide, "\n".join(offerings), 1.5, 2.7, 7.0, 3.2,
                font_size=13, color=WHITE)

    cta_box = slide.shapes.add_shape(1, Inches(2.5), Inches(5.8), Inches(5.0), Inches(0.8))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = ORANGE
    cta_box.line.fill.background()

    tb = cta_box.text_frame
    tb.text = "Schedule Your Strategy Call → monarchanalytics.in"
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(13)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK_BG


def slide_methodology(prs):
    """Slide 18: Data Sources & Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Data Sources & Methodology",
                       "How this audit was conducted")

    content = (
        "Data Sources Analyzed:\n"
        "  • Sales CRM Data (transactions, deal stages, pipeline)\n"
        "  • Marketing Spend Data (channels, campaigns, ROAS)\n"
        "  • Customer Database (LTV, retention, segments)\n"
        "  • Product Catalog (pricing, margins, cost structure)\n"
        "  • Operational Metrics (departmental KPIs)\n\n"
        "Methodology:\n"
        "  • Standardized data intake with format validation\n"
        "  • Automated data cleaning and normalization\n"
        "  • Rule-based revenue leak detection engine\n"
        "  • Industry benchmark comparison\n"
        "  • Pareto analysis for revenue concentration\n"
        "  • Cohort-based customer segmentation"
    )
    add_textbox(slide, content, 0.5, 1.4, 9.0, 5.0, font_size=13, color=WHITE)


def slide_about(prs):
    """Slide 19: About Monarch Analytics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "About Monarch Analytics",
                       "Your Business Intelligence Partner")

    about = (
        "Monarch Analytics is a data intelligence firm that helps growth-stage\n"
        "businesses unlock revenue trapped in underperforming systems.\n\n"
        "We productize every audit — so your insights are delivered fast,\n"
        "consistently, and with measurable financial impact.\n\n"
        "Our Approach:\n"
        "  • Standardized audit framework across all client engagements\n"
        "  • Automated KPI calculation and revenue leak detection\n"
        "  • Interactive dashboards and automated reporting\n"
        "  • Actionable recommendations backed by data\n\n"
        "Contact: hello@monarchanalytics.in  |  monarchanalytics.in"
    )
    add_textbox(slide, about, 0.5, 1.4, 9.0, 5.0, font_size=13, color=WHITE)


def slide_thank_you(prs):
    """Slide 20: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(slide, "MONARCH ANALYTICS", 0.4, 0.12, 4, 0.6,
                font_size=13, bold=True, color=ORANGE)

    add_textbox(slide, "Thank You", 1.0, 2.0, 8.0, 1.2,
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "This report was automatically generated by the\nMonarch Analytics Audit System.",
                1.0, 3.4, 8.0, 0.9, font_size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Data-Driven Decisions. Revenue Growth. Business Clarity.",
                1.0, 4.5, 8.0, 0.6, font_size=13, color=ORANGE, align=PP_ALIGN.CENTER)

    footer_bar = slide.shapes.add_shape(1, Inches(0), Inches(6.2), Inches(10), Inches(0.8))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = BLUE
    footer_bar.line.fill.background()

    add_textbox(slide, "CONFIDENTIAL  |  © Monarch Analytics  |  monarchanalytics.in",
                0.4, 6.28, 9.2, 0.5, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main Generator ────────────────────────────────────────────────────────────

def generate_report(kpis=None, leaks=None, drivers=None):
    """Build and save the full PowerPoint report."""
    print("\n📊 Generating PowerPoint Audit Report...")

    if kpis is None:
        kpis = load_kpis()
    if leaks is None:
        leaks = load_leaks()

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7)

    os.makedirs("audit_workspace/report_outputs", exist_ok=True)

    slides_fn = [
        lambda: slide_title(prs, kpis),
        lambda: slide_executive_summary(prs, kpis, leaks),
        lambda: slide_key_metrics(prs, kpis),
        lambda: slide_revenue_chart(prs),
        lambda: slide_sales_funnel(prs),
        lambda: slide_marketing(prs),
        lambda: slide_cac_ltv(prs),
        lambda: slide_marketing_funnel(prs),
        lambda: slide_customer_ltv(prs),
        lambda: slide_customer_segments(prs),
        lambda: slide_pareto(prs, drivers),
        lambda: slide_revenue_leaks(prs, leaks),
        lambda: slide_fix_improve_ignore(prs, leaks),
        lambda: slide_recommendations(prs, leaks),
        lambda: slide_roadmap(prs),
        lambda: slide_financial_impact(prs, leaks, kpis),
        lambda: slide_methodology(prs),
        lambda: slide_about(prs),
        lambda: slide_retainer_offer(prs),
        lambda: slide_thank_you(prs),
    ]

    for i, fn in enumerate(slides_fn, 1):
        fn()
        print(f"   🖼️  Slide {i:02d}/20 built")

    prs.save(OUT_PATH)
    print(f"\n   ✅ Report saved: {OUT_PATH}")
    return OUT_PATH


if __name__ == "__main__":
    generate_report()
